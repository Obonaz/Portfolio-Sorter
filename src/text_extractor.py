import os
from docx import Document
from pypdf import PdfReader
from openpyxl import load_workbook
from pptx import Presentation

# Custom Exception
class TextExtractionError(Exception):
    """Custom exception for errors during text extraction."""
    pass

# Supported file extension groups
DOC_EXTENSIONS = ('.doc', '.docx')
PDF_EXTENSIONS = ('.pdf',)
EXCEL_EXTENSIONS = ('.xls', '.xlsx')
POWERPOINT_EXTENSIONS = ('.ppt', '.pptx')
ALL_SUPPORTED_EXTENSIONS = DOC_EXTENSIONS + PDF_EXTENSIONS + EXCEL_EXTENSIONS + POWERPOINT_EXTENSIONS

# Helper functions
def _extract_text_from_docx(file_path: str) -> str:
    """Extracts text from a .docx file."""
    try:
        doc = Document(file_path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return '\n'.join(full_text)
    except Exception as e:
        raise TextExtractionError(f"Failed to extract text from DOCX '{file_path}': {e}")

def _extract_text_from_pdf(file_path: str) -> str:
    """Extracts text from a .pdf file."""
    try:
        reader = PdfReader(file_path)
        full_text = []
        for page in reader.pages:
            full_text.append(page.extract_text() or "") # Ensure None is handled
        return '\n'.join(full_text)
    except Exception as e:
        raise TextExtractionError(f"Failed to extract text from PDF '{file_path}': {e}")

def _extract_text_from_excel(file_path: str) -> str:
    """Extracts text from an .xls or .xlsx file."""
    try:
        workbook = load_workbook(filename=file_path, read_only=True, data_only=True)
        full_text = []
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value is not None:
                        full_text.append(str(cell.value))
        return '\n'.join(full_text)
    except Exception as e:
        raise TextExtractionError(f"Failed to extract text from Excel '{file_path}': {e}")

def _extract_text_from_pptx(file_path: str) -> str:
    """Extracts text from a .ppt or .pptx file."""
    try:
        prs = Presentation(file_path)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        return '\n'.join(full_text)
    except Exception as e:
        raise TextExtractionError(f"Failed to extract text from PowerPoint '{file_path}': {e}")

# Main extraction function
def extract_text(file_path: str) -> str | None:
    """
    Extracts text from a supported document file.

    Args:
        file_path: The path to the document file.

    Returns:
        The extracted text as a string, or None if the file type is unsupported.

    Raises:
        TextExtractionError: If text extraction fails for a supported file type.
        FileNotFoundError: If the file_path does not exist.
    """
    if not os.path.isfile(file_path):
        raise FileNotFoundError(f"File not found: '{file_path}'")

    file_extension = os.path.splitext(file_path)[1].lower()

    try:
        if file_extension in DOC_EXTENSIONS:
            # .doc files are not supported by python-docx directly,
            # this will only work for .docx
            if file_extension == '.doc':
                 print(f"Warning: Basic DOCX extractor cannot process legacy '.doc' files: {file_path}. Try converting to .docx.")
                 return None # Or raise specific error/warning
            return _extract_text_from_docx(file_path)
        elif file_extension in PDF_EXTENSIONS:
            return _extract_text_from_pdf(file_path)
        elif file_extension in EXCEL_EXTENSIONS:
            return _extract_text_from_excel(file_path)
        elif file_extension in POWERPOINT_EXTENSIONS:
            return _extract_text_from_pptx(file_path)
        elif file_extension not in ALL_SUPPORTED_EXTENSIONS:
            print(f"Unsupported file type for text extraction: '{file_path}' (extension: {file_extension})")
            return None
        else:
            # This case should ideally not be reached if ALL_SUPPORTED_EXTENSIONS is comprehensive
            print(f"File type '{file_extension}' is listed in a specific group but not ALL_SUPPORTED_EXTENSIONS for '{file_path}'. This is a bug.")
            return None
    except TextExtractionError: # Re-raise specific errors
        raise
    except Exception as e: # Catch-all for unexpected issues in dispatch logic
        raise TextExtractionError(f"An unexpected error occurred while processing '{file_path}': {e}")
